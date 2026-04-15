"""Roundtrip: MD → PPTX → MD with _class preservation via slide notes."""
import tempfile
from pathlib import Path

from marp_pptx.parser import parse_marp
from marp_pptx.theme import ThemeConfig, get_default_theme_path
from marp_pptx.builder import PptxBuilder
from marp_pptx.pptx2md import pptx_to_md_with_report, extract_slide

from pptx import Presentation


def _build_to_pptx(md_text: str, tmp_path: Path) -> Path:
    md_path = tmp_path / "input.md"
    md_path.write_text(md_text, encoding="utf-8")
    tc = ThemeConfig.from_css(get_default_theme_path())
    slides = parse_marp(str(md_path))
    b = PptxBuilder(base_path=tmp_path, theme=tc)
    b.build_all(slides)
    out = tmp_path / "out.pptx"
    b.save(str(out))
    return out


def test_builder_writes_class_to_notes(tmp_path: Path):
    md = """---
marp: true
---

<!-- _class: title -->
# Hello
## Sub

---

<!-- _class: kpi -->
# Metrics
<div class="kpi-container">
<div><span class="kpi-value">98%</span><span class="kpi-label">Acc</span></div>
</div>

---

<!-- _class: end -->
# Thank You
"""
    pptx = _build_to_pptx(md, tmp_path)
    prs = Presentation(str(pptx))
    slides = list(prs.slides)
    assert len(slides) >= 3
    notes = [s.notes_slide.notes_text_frame.text for s in slides]
    assert "_class: title" in notes[0]
    assert "_class: kpi" in notes[1]
    assert "_class: end" in notes[2]


def test_pptx2md_recovers_class_from_notes(tmp_path: Path):
    md = """---
marp: true
---

<!-- _class: title -->
# A

---

<!-- _class: funnel -->
# Funnel
<div class="fn-container">
<div><span class="fn-label">Apply</span><span class="fn-value">1000</span></div>
<div><span class="fn-label">Hired</span><span class="fn-value">10</span></div>
</div>

---

<!-- _class: pros-cons -->
# Compare
<div class="pc-pros">
<ul><li>Fast</li></ul>
</div>
<div class="pc-cons">
<ul><li>Costly</li></ul>
</div>
"""
    pptx = _build_to_pptx(md, tmp_path)
    report = pptx_to_md_with_report(pptx)
    cls_list = [s["inferred_class"] for s in report["slides"]]
    sources = [s["source"] for s in report["slides"]]
    assert "title" in cls_list
    assert "funnel" in cls_list
    assert "pros-cons" in cls_list
    # All should come from notes (our builder-created PPTXs)
    assert all(src == "notes" for src in sources), f"expected all notes, got {sources}"


def test_math_mode_png_produces_picture_not_omml(tmp_path: Path):
    """theme.math_mode='png' forces matplotlib PNG fallback instead of OMML."""
    import shutil
    md = """---
marp: true
---

<!-- _class: equation -->
# F
<div class="eq-main">
$$E = mc^2$$
</div>
"""
    # png mode
    md_path = tmp_path / "in.md"
    md_path.write_text(md, encoding="utf-8")
    tc = ThemeConfig.from_css(get_default_theme_path())
    tc.math_mode = "png"
    slides = parse_marp(str(md_path))
    b = PptxBuilder(base_path=tmp_path, theme=tc)
    b.build_all(slides)
    out = tmp_path / "png.pptx"
    b.save(str(out))

    # Inspect: should contain a picture shape and NO <a14:m> OMML element
    import zipfile
    with zipfile.ZipFile(str(out)) as z:
        slide2 = z.read("ppt/slides/slide1.xml").decode(errors="ignore")
    assert "a14:m" not in slide2, "OMML should not be emitted in png math_mode"
    # Picture shape indicator: <p:pic>
    assert "p:pic" in slide2, "expected a picture shape for math when math_mode=png"


def test_math_mode_omml_is_default(tmp_path: Path):
    """Default math_mode='omml' produces native OMML when Pandoc is available."""
    import shutil
    if shutil.which("pandoc") is None:
        return  # skip — OMML requires pandoc
    md = """---
marp: true
---

<!-- _class: equation -->
# F
<div class="eq-main">
$$E = mc^2$$
</div>
"""
    md_path = tmp_path / "in.md"
    md_path.write_text(md, encoding="utf-8")
    tc = ThemeConfig.from_css(get_default_theme_path())
    # Default math_mode should be 'omml'
    assert tc.math_mode == "omml"
    slides = parse_marp(str(md_path))
    b = PptxBuilder(base_path=tmp_path, theme=tc)
    b.build_all(slides)
    out = tmp_path / "omml.pptx"
    b.save(str(out))
    import zipfile
    with zipfile.ZipFile(str(out)) as z:
        slide = z.read("ppt/slides/slide1.xml").decode(errors="ignore")
    assert "a14:m" in slide, "expected native OMML element for default math_mode"


def test_pptx2md_default_no_class_note_means_heuristic(tmp_path: Path):
    """Slides without an explicit _class should have note written as 'default'
    and pptx2md should recognize this as a null class (fall back to heuristic).
    """
    md = """---
marp: true
---

# Plain slide title
- bullet
"""
    pptx = _build_to_pptx(md, tmp_path)
    prs = Presentation(str(pptx))
    s0 = list(prs.slides)[0]
    note_text = s0.notes_slide.notes_text_frame.text
    assert "_class: default" in note_text
    # extract_slide should yield note_class = None (because default maps to None)
    es = extract_slide(0, s0)
    assert es.note_class is None
