"""PptxBuilder: converts parsed SlideData into PowerPoint slides.

All theme-dependent values come from self.theme (ThemeConfig instance),
eliminating the global state of the original convert_v2.py.
"""
from __future__ import annotations

import re
import sys
import tempfile
import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from marp_pptx.parser import SlideData, strip_html
from marp_pptx.theme import ThemeConfig
from marp_pptx.layout import (
    SW, SH, MARGIN_L, MARGIN_R, MARGIN_T, CONTENT_W,
    TITLE_H, TITLE_TOP, BODY_TOP, BODY_H,
    SZ_TITLE, SZ_H2, SZ_H3, SZ_BODY, SZ_COL, SZ_SMALL,
    SZ_FOOT, SZ_EQ, SZ_EQ_VAR, SZ_ZONE_L, SZ_ZONE_B,
)
from marp_pptx.math.omml import latex_to_omml_element, OmmlError
from marp_pptx.math.renderer import render_latex_png

try:
    import cairosvg
    HAS_CAIROSVG = True
except ImportError:
    HAS_CAIROSVG = False

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


class PptxBuilder:
    def __init__(self, base_path: Path, theme: ThemeConfig):
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self.base_path = base_path
        self.theme = theme
        self._img_cache: dict = {}

    def _fs(self, pt_val):
        """Scale a Pt value by theme.font_scale. Min 8pt.

        Accepts either a Pt object or a raw int (points).
        """
        from pptx.util import Pt as _Pt
        scale = getattr(self.theme, "font_scale", 1.0)
        try:
            base = pt_val.pt
        except AttributeError:
            base = float(pt_val)
        return _Pt(max(8, base * scale))

    def save(self, path: str):
        self._ensure_ea_font()
        self.prs.save(path)

    # ── Theme shortcuts ──
    @property
    def PRIMARY(self): return self.theme.primary
    @property
    def SECONDARY(self): return self.theme.secondary
    @property
    def ACCENT(self): return self.theme.accent
    @property
    def FG(self): return self.theme.fg
    @property
    def MUTED(self): return self.theme.muted
    @property
    def LIGHT(self): return self.theme.light
    @property
    def WHITE(self): return self.theme.white
    @property
    def FONT(self): return self.theme.font
    @property
    def FONT_HEAD(self): return self.theme.font_head
    @property
    def FONT_EA(self): return self.theme.font_ea
    @property
    def FONT_MONO(self): return self.theme.font_mono
    @property
    def LAYOUT(self): return self.theme.layout

    # ── EA font injection ──
    def _ensure_ea_font(self):
        rpr_tags = (f"{{{NS_A}}}rPr", f"{{{NS_A}}}defRPr", f"{{{NS_A}}}endParaRPr")
        for slide in self.prs.slides:
            root = slide._element
            for tag in rpr_tags:
                for rpr in root.iter(tag):
                    self._patch_rpr(rpr)

    def _patch_rpr(self, rpr):
        latin = rpr.find(f"{{{NS_A}}}latin")
        if latin is None:
            return
        ea = rpr.find(f"{{{NS_A}}}ea")
        if ea is None:
            ea = etree.Element(f"{{{NS_A}}}ea")
            ea.set("typeface", self.FONT_EA)
            latin.addnext(ea)
        else:
            ea.set("typeface", self.FONT_EA)
        cs = rpr.find(f"{{{NS_A}}}cs")
        if cs is None:
            cs = etree.Element(f"{{{NS_A}}}cs")
            cs.set("typeface", self.FONT_EA)
            ea.addnext(cs)
        else:
            cs.set("typeface", self.FONT_EA)

    # ── Math helpers ──
    def _omml_element(self, latex: str, display: bool):
        try:
            return latex_to_omml_element(latex, display=display)
        except OmmlError as e:
            print(f"  OMML failed: {latex[:40]}... ({e}) — falling back to PNG", file=sys.stderr)
            return None

    def _add_math_omml_display(self, slide, latex, left, top, width, pt_size=28):
        el = self._omml_element(latex, display=True)
        if el is None:
            return None
        height = Pt(int(pt_size * 1.7))
        tb = self._add_textbox(slide, left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = ""
        run.font.name = self.FONT
        run.font.size = self._fs(Pt(pt_size))
        run.font.color.rgb = self.FG
        p._p.append(el)
        return tb

    def _append_math_omml_inline(self, para, latex, size, color):
        el = self._omml_element(latex, display=False)
        if el is None:
            return False
        run = para.add_run()
        run.text = ""
        run.font.name = self.FONT
        run.font.size = size
        if color is not None:
            run.font.color.rgb = color
        para._p.append(el)
        return True

    def _render_math(self, latex: str, display: bool = False, fontsize: int = 28) -> str | None:
        """Render LaTeX to PNG via matplotlib. Returns path to PNG."""
        color_hex = f"#{self.FG}"
        return render_latex_png(latex, fontsize=fontsize, display=display, color=color_hex)

    def _add_math_image(self, slide, latex, left, top, max_width, display=True, fontsize=28):
        png = self._render_math(latex, display=display, fontsize=fontsize)
        if not png:
            return None
        from PIL import Image
        with Image.open(png) as im:
            iw, ih = im.size
        dpi = 150
        pw = int(iw * 914400 / dpi)
        ph = int(ih * 914400 / dpi)
        if pw > max_width:
            scale = max_width / pw
            pw = int(pw * scale)
            ph = int(ph * scale)
        img_left = left + (max_width - pw) // 2
        slide.shapes.add_picture(png, img_left, top, pw, ph)
        return (pw, ph)

    # ── Basic shape helpers ──
    def _blank_slide(self):
        layout = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(layout)

    def _add_textbox(self, slide, left, top, width, height):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        return tb

    def _set_bg(self, slide, color):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = color

    def _set_gradient_bg(self, slide, c1, c2):
        bg = slide.background
        bg.fill.gradient()
        bg.fill.gradient_stops[0].color.rgb = c1
        bg.fill.gradient_stops[0].position = 0.0
        bg.fill.gradient_stops[1].color.rgb = c2
        bg.fill.gradient_stops[1].position = 1.0

    def _add_title(self, slide, text, top=None, color=None):
        if color is None:
            color = self.PRIMARY
        if top is None:
            top = TITLE_TOP
        deco_color_map = {"primary": self.PRIMARY, "secondary": self.SECONDARY, "accent": self.ACCENT}
        deco_c = deco_color_map.get(self.LAYOUT.h1_deco_color, self.PRIMARY)
        deco_w = Pt(self.LAYOUT.h1_deco_width)
        text_left = MARGIN_L
        text_w = CONTENT_W

        if self.LAYOUT.h1_deco == "left-bar":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(MARGIN_L), int(top), int(deco_w), int(TITLE_H))
            bar.fill.solid()
            bar.fill.fore_color.rgb = deco_c
            bar.line.fill.background()
            text_left = int(MARGIN_L + deco_w + Pt(10))
            text_w = int(CONTENT_W - deco_w - Pt(10))
        elif self.LAYOUT.h1_deco == "bottom-line":
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(MARGIN_L), int(top + TITLE_H - Pt(2)),
                int(CONTENT_W), int(deco_w))
            line.fill.solid()
            line.fill.fore_color.rgb = deco_c
            line.line.fill.background()
        elif self.LAYOUT.h1_deco == "top-line":
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(MARGIN_L), int(top), int(CONTENT_W), int(deco_w))
            line.fill.solid()
            line.fill.fore_color.rgb = deco_c
            line.line.fill.background()
            top = int(top + deco_w + Pt(4))
        elif self.LAYOUT.h1_deco == "double-bottom":
            for offset in [0, Pt(6)]:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, int(MARGIN_L),
                    int(top + TITLE_H - Pt(2) + offset),
                    int(CONTENT_W), Pt(2))
                line.fill.solid()
                line.fill.fore_color.rgb = deco_c
                line.line.fill.background()

        tb = self._add_textbox(slide, int(text_left), int(top), int(text_w), TITLE_H)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.FONT_HEAD
        p.font.size = self._fs(SZ_TITLE)
        p.font.bold = True
        p.font.color.rgb = color
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        return tb

    def _add_para(self, tf, text, size=None, color=None, bold=False, italic=False, space_before=Pt(4)):
        if size is None:
            size = SZ_BODY
        if color is None:
            color = self.FG
        p = tf.add_paragraph()
        p.text = text
        p.font.name = self.FONT
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_before = space_before
        return p

    def _add_body_text(self, slide, lines, left=None, top=None, width=None, height=None, size=None):
        if size is None:
            size = SZ_BODY
        if left is None: left = MARGIN_L
        if top is None: top = BODY_TOP
        if width is None: width = CONTENT_W
        if height is None: height = BODY_H

        tb = self._add_textbox(slide, left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True

        first = True
        for line in lines:
            s = line.strip()
            if not s:
                continue
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            is_h2 = s.startswith("## ")
            is_h3 = s.startswith("### ")
            is_bullet = s.startswith("- ") or s.startswith("* ")
            is_numbered = re.match(r"^\d+\.\s", s)

            if is_h2:
                p.text = strip_html(s[3:])
                p.font.name = self.FONT_HEAD
                p.font.size = self._fs(SZ_H2)
                p.font.bold = True
                p.font.color.rgb = self.SECONDARY
                p.space_before = Pt(10)
            elif is_h3:
                p.text = strip_html(s[4:])
                p.font.name = self.FONT_HEAD
                p.font.size = self._fs(SZ_H3)
                p.font.bold = True
                p.font.color.rgb = self.MUTED
                p.space_before = Pt(6)
            elif is_bullet:
                # Apply rich text (bold markdown **...**) to bullet content
                self._set_rich_text(p, s[2:], size, self.FG)
                p.level = 0
                p.space_before = Pt(4)
                pPr = p._p.get_or_add_pPr()
                buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
                for existing in pPr.findall(qn("a:buChar")):
                    pPr.remove(existing)
                for existing in pPr.findall(qn("a:buNone")):
                    pPr.remove(existing)
                pPr.append(buChar)
            elif is_numbered:
                p.text = s
                p.font.name = self.FONT
                p.font.size = size
                p.font.color.rgb = self.FG
                p.space_before = Pt(4)
            else:
                self._set_rich_text(p, s, size, self.FG)
                p.space_before = Pt(4)

        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        return tb

    def _add_plain_run(self, para, text, size, color, bold=False, mono=False):
        """Append a single styled run to para. No-op if text is empty."""
        if not text:
            return
        run = para.add_run()
        run.text = text
        run.font.name = self.FONT_MONO if mono else self.FONT
        run.font.size = self._fs(size) if hasattr(size, "pt") else size
        if color is not None:
            run.font.color.rgb = color
        if bold:
            run.font.bold = True

    # Combined inline markup: **bold**, `code`, $math$
    _RICH_PATTERN = re.compile(
        r"(\*\*[^\*\n]+?\*\*)"
        r"|(`[^`\n]+?`)"
        r"|(\$[^\$\n]+?\$)"
    )

    def _set_rich_text(self, para, text, size=None, color=None):
        """Render inline markup in a SINGLE paragraph / SINGLE textbox.

        Handles **bold**, `code` (monospace), and $math$ (OMML) without
        breaking the containing textbox. Runs co-exist with Japanese+Latin
        mixed text via the ea-font patch applied at save time.
        """
        if size is None:
            size = SZ_BODY
        if color is None:
            color = self.FG
        para.clear()
        if not text:
            return

        pos = 0
        for m in self._RICH_PATTERN.finditer(text):
            if m.start() > pos:
                self._add_plain_run(para, text[pos:m.start()], size, color)
            if m.group(1):  # **bold**
                self._add_plain_run(para, m.group(1)[2:-2], size, color, bold=True)
            elif m.group(2):  # `code`
                self._add_plain_run(para, m.group(2)[1:-1], size, color, mono=True)
            elif m.group(3):  # $math$
                latex = m.group(3)[1:-1]
                if not self._append_math_omml_inline(para, latex, size, color):
                    self._add_plain_run(para, m.group(3), size, color)
            pos = m.end()
        if pos < len(text):
            self._add_plain_run(para, text[pos:], size, color)

    # Backward-compatible alias for callers that used the math-only name
    def _set_text_with_inline_math(self, para, text, size, color):
        return self._set_rich_text(para, text, size=size, color=color)

    def _resolve_image(self, img_path: str) -> str | None:
        p = self.base_path / img_path
        if not p.exists():
            return None
        if p.suffix.lower() == ".svg":
            png_path = p.with_suffix(".png")
            if not png_path.exists() or png_path.stat().st_mtime < p.stat().st_mtime:
                if HAS_CAIROSVG:
                    cairosvg.svg2png(url=str(p), write_to=str(png_path), output_width=1400, dpi=300)
                else:
                    return None
            return str(png_path)
        return str(p)

    def _fill_multiline_box(self, tf, text, size, color):
        """Fill a textbox with text that may contain bullets and continuation lines.

        Handles:
        - `- item` and `* item` as bullets
        - Continuation lines indented MORE than their parent bullet
        - **bold** markdown via _set_rich_text
        """
        # Strip common leading whitespace (dedent) so HTML-source indentation
        # doesn't get mistaken for content continuation.
        raw = [line for line in text.split("\n") if line.strip()]
        if not raw:
            return
        indents = [len(l) - len(l.lstrip()) for l in raw]
        base_indent = min(indents) if indents else 0
        dedented = [l[base_indent:] if len(l) >= base_indent else l for l in raw]

        # Merge continuation: a line is a continuation only if the previous
        # line was a bullet AND this line is further indented.
        merged: list[str] = []
        last_was_bullet = False
        for line in dedented:
            stripped = line.lstrip()
            line_indent = len(line) - len(stripped)
            is_bullet = stripped.startswith("- ") or stripped.startswith("* ")
            if last_was_bullet and line_indent > 0 and not is_bullet and merged:
                merged[-1] = merged[-1] + " " + stripped
            else:
                merged.append(stripped)
                last_was_bullet = is_bullet

        first = True
        for line in merged:
            s = line.strip()
            if not s:
                continue
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            is_bullet = s.startswith("- ") or s.startswith("* ")
            if is_bullet:
                self._set_rich_text(p, s[2:], size, color)
                p.space_before = Pt(4)
                pPr = p._p.get_or_add_pPr()
                buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
                for existing in pPr.findall(qn("a:buChar")):
                    pPr.remove(existing)
                pPr.append(buChar)
            else:
                self._set_rich_text(p, s, size, color)
                p.space_before = Pt(4)

    def _add_accent_box(self, slide, text, left, top, width, height, border_color=None):
        if border_color is None:
            border_color = self.ACCENT
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.LIGHT
        bg.line.fill.background()
        bdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(6), height)
        bdr.fill.solid()
        bdr.fill.fore_color.rgb = border_color
        bdr.line.fill.background()
        tb = self._add_textbox(slide, left + Pt(16), top + Pt(8), width - Pt(32), height - Pt(16))
        tf = tb.text_frame
        tf.word_wrap = True
        self._fill_multiline_box(tf, text, SZ_BODY, self.FG)
        return tb

    def _add_conclusion_box(self, slide, text, left, top, width, height):
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        bg.adjustments[0] = 0.02
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.LIGHT
        bg.line.fill.background()
        tb = self._add_textbox(slide, left + Pt(16), top + Pt(10), width - Pt(32), height - Pt(20))
        tf = tb.text_frame
        tf.word_wrap = True
        self._fill_multiline_box(tf, text, SZ_BODY, self.FG)
        return tb

    def _add_footnote(self, slide, text):
        left = MARGIN_L
        top = SH - Inches(0.55)
        width = CONTENT_W
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
        ln.fill.solid()
        ln.fill.fore_color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
        ln.line.fill.background()
        tb = self._add_textbox(slide, left, top + Pt(4), width, Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.FONT
        p.font.size = self._fs(SZ_FOOT)
        p.font.color.rgb = self.MUTED
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    def _add_zone_box(self, slide, left, top, width, height,
                      label="", body="", fill_color=None,
                      label_size=None, body_size=None):
        if fill_color is None:
            fill_color = self.LIGHT
        if label_size is None:
            label_size = SZ_ZONE_L
        if body_size is None:
            body_size = SZ_ZONE_B
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        bg.adjustments[0] = self.LAYOUT.box_radius
        if self.LAYOUT.box_style in ("filled", "card"):
            bg.fill.solid()
            bg.fill.fore_color.rgb = fill_color
            if self.LAYOUT.box_style == "card":
                bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                bg.line.width = Pt(1)
            else:
                bg.line.fill.background()
        elif self.LAYOUT.box_style == "accent-border":
            bg.fill.background()
            bg.line.color.rgb = self.ACCENT
            bg.line.width = Pt(1.5)
        else:
            bg.fill.background()
            bg.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
            bg.line.width = Pt(0.75)
        pad = Pt(14)
        tb = self._add_textbox(slide, left + pad, top + pad, width - pad * 2, height - pad * 2)
        tf = tb.text_frame
        tf.word_wrap = True
        if label:
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = label
            run.font.name = self.FONT_HEAD
            run.font.size = label_size
            run.font.bold = True
            run.font.color.rgb = self.SECONDARY
        if body:
            p2 = tf.add_paragraph() if label else tf.paragraphs[0]
            p2.space_before = Pt(6)
            self._set_text_with_inline_math(p2, body, body_size, self.FG)
        return bg, tb

    def _add_column_content(self, slide, lines, left, top, width, height, size=None):
        if size is None:
            size = SZ_COL
        text_lines = []
        images = []
        for line in lines:
            img_m = re.match(r"!\[(?:w:\d+)?\]\(([^)]+)\)", line.strip())
            if img_m:
                images.append(img_m.group(1))
            else:
                text_lines.append(line)
        cur_top = top
        for img_path in images:
            img_file = self._resolve_image(img_path)
            if img_file:
                from PIL import Image
                with Image.open(img_file) as im:
                    iw, ih = im.size
                max_w = int(width * 0.95)
                max_h = int(Inches(2.5))
                scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
                pw = int(iw * scale * 914400 / 96)
                ph = int(ih * scale * 914400 / 96)
                img_left = left + (width - pw) // 2
                slide.shapes.add_picture(img_file, img_left, cur_top, pw, ph)
                cur_top += ph + Inches(0.1)
        remaining_h = top + height - cur_top
        if text_lines and remaining_h > 0:
            self._add_body_text(slide, text_lines, left=left, top=int(cur_top),
                               width=int(width), height=int(remaining_h), size=size)

    # ══════════════════════════════════════════════
    # Slide type builders
    # ══════════════════════════════════════════════

    def build_title(self, sd: SlideData):
        slide = self._blank_slide()
        if self.LAYOUT.title_bg == "gradient":
            self._set_gradient_bg(slide, self.PRIMARY, self.SECONDARY)
        elif self.LAYOUT.title_bg == "dark":
            self._set_bg(slide, self.PRIMARY)
        elif self.LAYOUT.title_bg == "light":
            self._set_bg(slide, self.LIGHT)
        is_dark = self.LAYOUT.title_bg in ("gradient", "dark")
        align = PP_ALIGN.CENTER if self.LAYOUT.title_align == "center" else PP_ALIGN.LEFT
        h_color = self.WHITE if is_dark else self.PRIMARY
        sub_color = RGBColor(0xCC, 0xCC, 0xCC) if is_dark else self.MUTED
        tb = self._add_textbox(slide, Inches(1), Inches(1.5), SW - Inches(2), Inches(2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sd.h1
        p.font.name = self.FONT_HEAD
        p.font.size = self._fs(Pt(44))
        p.font.bold = True
        p.font.color.rgb = h_color
        p.alignment = align
        tb2 = self._add_textbox(slide, Inches(1), Inches(3.2), SW - Inches(2), Inches(3.5))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        remaining = []
        past_h1 = False
        for line in sd.raw.split("\n"):
            s = line.strip()
            if s.startswith("# ") and not past_h1:
                past_h1 = True
                continue
            if s.startswith("## "):
                remaining.append(s[3:])
                continue
            if past_h1 and s:
                remaining.append(strip_html(s))
        first = True
        for line in remaining:
            if first:
                p = tf2.paragraphs[0]
                first = False
            else:
                p = tf2.add_paragraph()
            p.text = line
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(20))
            p.font.color.rgb = sub_color
            p.alignment = align
            p.space_before = Pt(6)

    def build_divider(self, sd: SlideData):
        slide = self._blank_slide()
        align = PP_ALIGN.CENTER if self.LAYOUT.divider_align == "center" else PP_ALIGN.LEFT
        x = Inches(1) if self.LAYOUT.divider_align == "center" else Inches(1.5)
        w = SW - Inches(2) if self.LAYOUT.divider_align == "center" else SW - Inches(3)
        tb = self._add_textbox(slide, x, Inches(2.5), w, Inches(1.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = sd.h1
        p.font.name = self.FONT_HEAD
        p.font.size = self._fs(Pt(36))
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY
        p.alignment = align
        if sd.h2:
            p2 = tf.add_paragraph()
            p2.text = sd.h2
            p2.font.name = self.FONT
            p2.font.size = self._fs(Pt(22))
            p2.font.color.rgb = self.MUTED
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(12)

    def build_default(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        cur_top = BODY_TOP
        if sd.body_lines:
            tb = self._add_body_text(slide, sd.body_lines, top=int(cur_top))
            # Advance cur_top by the body height if there's also a table
            if sd.table_rows:
                cur_top = BODY_TOP + Inches(0.5) * len([l for l in sd.body_lines if l.strip()])
        if sd.table_rows:
            rows = len(sd.table_rows)
            cols = max(len(r) for r in sd.table_rows) if sd.table_rows else 1
            tbl_h = min(Inches(4.5), Inches(0.45) * rows)
            table = slide.shapes.add_table(rows, cols, MARGIN_L, int(cur_top), CONTENT_W, tbl_h).table
            for ri, row in enumerate(sd.table_rows):
                for ci, cell_text in enumerate(row):
                    if ci >= cols:
                        break
                    cell = table.cell(ri, ci)
                    cell.text = strip_html(cell_text)
                    for p in cell.text_frame.paragraphs:
                        p.font.name = self.FONT
                        p.font.size = self._fs(SZ_SMALL)
                        p.font.color.rgb = self.WHITE if ri == 0 else self.FG
                    if ri == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.PRIMARY
                    elif ri % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.LIGHT
        if sd.bottom_text:
            self._add_accent_box(slide, sd.bottom_text, MARGIN_L, SH - Inches(1.8), CONTENT_W, Inches(1.0))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_equation(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        eq_top = BODY_TOP + Inches(0.2)
        omml_box = self._add_math_omml_display(slide, sd.eq_main, MARGIN_L, eq_top, CONTENT_W, pt_size=28)
        if omml_box is not None:
            var_top = eq_top + omml_box.height + Inches(0.25)
        else:
            result = self._add_math_image(slide, sd.eq_main, MARGIN_L, eq_top, CONTENT_W, display=True, fontsize=36)
            if result:
                _, eq_h = result
                var_top = eq_top + eq_h + Inches(0.3)
            else:
                tb = self._add_textbox(slide, MARGIN_L, eq_top, CONTENT_W, Inches(1.2))
                tf = tb.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = sd.eq_main
                p.font.name = self.FONT
                p.font.size = self._fs(Pt(28))
                p.font.color.rgb = self.FG
                p.alignment = PP_ALIGN.CENTER
                var_top = eq_top + Inches(1.6)
        if sd.eq_vars:
            desc_left = Inches(2.0)
            row_h = Inches(0.58)
            for vi, (sym, desc) in enumerate(sd.eq_vars):
                row_top = var_top + int(row_h * vi)
                sym_latex = sym.strip().strip("$")
                stb = self._add_textbox(slide, desc_left, row_top, Inches(2.0), row_h)
                stf = stb.text_frame
                stf.vertical_anchor = MSO_ANCHOR.MIDDLE
                sp = stf.paragraphs[0]
                sp.alignment = PP_ALIGN.RIGHT
                if not self._append_math_omml_inline(sp, sym_latex, Pt(22), self.SECONDARY):
                    sp.text = sym
                    sp.font.name = self.FONT
                    sp.font.size = self._fs(Pt(18))
                    sp.font.bold = True
                    sp.font.color.rgb = self.SECONDARY
                dtb = self._add_textbox(slide, desc_left + Inches(2.3), row_top, Inches(6.5), row_h)
                dtf = dtb.text_frame
                dtf.word_wrap = True
                dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
                self._set_text_with_inline_math(dtf.paragraphs[0], desc, Pt(17), self.FG)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_equations(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.eq_system)
        if n == 0:
            return
        top = BODY_TOP + Inches(0.1)
        label_left = MARGIN_L + Inches(0.3)
        label_w = Inches(1.9)
        eq_left = label_left + label_w + Inches(0.2)
        eq_w = CONTENT_W - (eq_left - MARGIN_L) - Inches(0.3)
        vars_h = Inches(0.58) * max(len(sd.eq_vars), 0)
        footnote_h = Inches(0.4) if sd.footnote else Inches(0)
        avail_h = BODY_H - vars_h - footnote_h - Inches(0.3)
        row_h = min(Inches(1.3), max(Inches(0.75), int(avail_h / n)))
        pt_size = 30 if n <= 3 else (26 if n <= 4 else 22)
        for i, (label, latex) in enumerate(sd.eq_system):
            row_top = top + int(row_h * i)
            if label:
                ltb = self._add_textbox(slide, label_left, row_top, label_w, row_h)
                ltf = ltb.text_frame
                ltf.vertical_anchor = MSO_ANCHOR.MIDDLE
                ltf.word_wrap = True
                lp = ltf.paragraphs[0]
                lp.text = label
                lp.font.name = self.FONT
                lp.font.size = self._fs(Pt(max(14, pt_size - 10)))
                lp.font.color.rgb = self.SECONDARY
                lp.alignment = PP_ALIGN.RIGHT
            el = self._omml_element(latex, display=True)
            etb = self._add_textbox(slide, eq_left, row_top, eq_w, row_h)
            etf = etb.text_frame
            etf.word_wrap = True
            etf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ep = etf.paragraphs[0]
            ep.alignment = PP_ALIGN.LEFT
            erun = ep.add_run()
            erun.text = ""
            erun.font.name = self.FONT
            erun.font.size = self._fs(Pt(pt_size))
            erun.font.color.rgb = self.FG
            if el is not None:
                ep._p.append(el)
            else:
                etb.element.getparent().remove(etb.element)
                result = self._add_math_image(slide, latex, eq_left, row_top, eq_w, display=True, fontsize=pt_size)
                if not result:
                    tb2 = self._add_textbox(slide, eq_left, row_top, eq_w, row_h)
                    tf2 = tb2.text_frame
                    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p2 = tf2.paragraphs[0]
                    p2.text = latex
                    p2.font.name = self.FONT
                    p2.font.size = self._fs(Pt(pt_size - 4))
                    p2.font.color.rgb = self.FG
        if sd.eq_vars:
            var_top = top + int(row_h * n) + Inches(0.25)
            desc_left = Inches(2.0)
            row_h_v = Inches(0.52)
            for vi, (sym, desc) in enumerate(sd.eq_vars):
                row_top = var_top + int(row_h_v * vi)
                sym_latex = sym.strip().strip("$")
                stb = self._add_textbox(slide, desc_left, row_top, Inches(2.0), row_h_v)
                stf = stb.text_frame
                stf.vertical_anchor = MSO_ANCHOR.MIDDLE
                sp = stf.paragraphs[0]
                sp.alignment = PP_ALIGN.RIGHT
                if not self._append_math_omml_inline(sp, sym_latex, Pt(18), self.SECONDARY):
                    sp.text = sym
                    sp.font.name = self.FONT
                    sp.font.size = self._fs(Pt(16))
                    sp.font.bold = True
                    sp.font.color.rgb = self.SECONDARY
                dtb = self._add_textbox(slide, desc_left + Inches(2.3), row_top, Inches(6.5), row_h_v)
                dtf = dtb.text_frame
                dtf.word_wrap = True
                dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
                self._set_text_with_inline_math(dtf.paragraphs[0], desc, Pt(15), self.FG)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_columns(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.columns)
        if n == 0:
            return
        gap = Inches(0.4)
        total_gap = gap * (n - 1)
        col_w = (CONTENT_W - total_gap) / n
        for i, col_lines in enumerate(sd.columns):
            left = MARGIN_L + i * (col_w + gap)
            self._add_column_content(slide, col_lines, left=int(left), top=BODY_TOP, width=int(col_w), height=BODY_H, size=Pt(18))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_sandwich(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        cur_top = BODY_TOP
        if sd.top_text:
            tb = self._add_textbox(slide, MARGIN_L, cur_top, CONTENT_W, Inches(0.8))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd.top_text
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(19))
            p.font.color.rgb = self.SECONDARY
            cur_top += Inches(1.0)
        n = len(sd.columns)
        if n > 0:
            gap = Inches(0.4)
            total_gap = gap * (n - 1)
            col_w = (CONTENT_W - total_gap) / n
            col_h = Inches(2.5)
            for i, col_lines in enumerate(sd.columns):
                left = MARGIN_L + i * (col_w + gap)
                self._add_body_text(slide, col_lines, left=left, top=int(cur_top), width=int(col_w), height=int(col_h), size=Pt(17))
            cur_top += col_h + Inches(0.2)
        if sd.bottom_text:
            remaining_h = SH - cur_top - Inches(0.4)
            box_h = min(Inches(1.0), int(remaining_h))
            self._add_conclusion_box(slide, sd.bottom_text, MARGIN_L, int(cur_top), CONTENT_W, box_h)

    def build_figure(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        img_file = self._resolve_image(sd.image_path) if sd.image_path else None
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(CONTENT_W * 0.85)
            max_h = int(Inches(3.5))
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96)
            ph = int(ih * scale * 914400 / 96)
            left = (SW - pw) // 2
            slide.shapes.add_picture(img_file, left, BODY_TOP, pw, ph)
            cap_top = BODY_TOP + ph + Inches(0.15)
        else:
            cap_top = BODY_TOP + Inches(3.5)
        if sd.caption:
            tb = self._add_textbox(slide, MARGIN_L, cap_top, CONTENT_W, Inches(0.8))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd.caption
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(14))
            p.font.color.rgb = self.FG
            p.alignment = PP_ALIGN.CENTER
        if sd.body_lines:
            desc_top = cap_top + Inches(0.7)
            self._add_body_text(slide, sd.body_lines, top=int(desc_top), size=Pt(17))

    def build_table(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        sub_top = BODY_TOP
        if sd.h2:
            tb = self._add_textbox(slide, MARGIN_L, sub_top, CONTENT_W, Inches(0.35))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.text = sd.h2
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(16))
            p.font.bold = True
            p.font.color.rgb = self.SECONDARY
            sub_top += Inches(0.5)
        if not sd.table_rows:
            return
        rows = len(sd.table_rows)
        cols = max(len(r) for r in sd.table_rows) if sd.table_rows else 1
        tbl_h = min(Inches(4.5), Inches(0.45) * rows)
        table = slide.shapes.add_table(rows, cols, MARGIN_L, sub_top, CONTENT_W, tbl_h).table
        for ri, row in enumerate(sd.table_rows):
            for ci, cell_text in enumerate(row):
                if ci >= cols:
                    break
                cell = table.cell(ri, ci)
                cell.text = strip_html(cell_text)
                for p in cell.text_frame.paragraphs:
                    p.font.name = self.FONT
                    p.font.size = self._fs(SZ_SMALL)
                    p.font.color.rgb = self.WHITE if ri == 0 else self.FG
                if ri == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.PRIMARY
                elif ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.LIGHT
        if sd.bottom_text:
            bt_top = sub_top + tbl_h + Inches(0.2)
            self._add_accent_box(slide, sd.bottom_text, MARGIN_L, bt_top, CONTENT_W, Inches(0.8))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_references(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.ref_items:
            return
        tb = self._add_textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, BODY_H)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, (author, title, venue) in enumerate(sd.ref_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"[{i+1}] "
            run.font.name = self.FONT
            run.font.size = self._fs(Pt(11))
            run.font.color.rgb = self.MUTED
            if author:
                run = p.add_run()
                run.text = author + " "
                run.font.name = self.FONT
                run.font.size = self._fs(Pt(11))
                run.font.bold = True
                run.font.color.rgb = self.FG
            if title:
                run = p.add_run()
                run.text = title + " "
                run.font.name = self.FONT
                run.font.size = self._fs(Pt(11))
                run.font.color.rgb = self.FG
            if venue:
                run = p.add_run()
                run.text = venue
                run.font.name = self.FONT
                run.font.size = self._fs(Pt(11))
                run.font.color.rgb = self.MUTED

    def build_timeline_h(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.timeline_items)
        if n == 0:
            return
        line_y = BODY_TOP + Inches(0.8)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, line_y, CONTENT_W, Pt(4))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x90, 0x9A, 0xA8)
        line.line.fill.background()
        item_w = CONTENT_W / n
        for i, item in enumerate(sd.timeline_items):
            cx = MARGIN_L + int(item_w * i) + int(item_w / 2)
            dot_color = self.ACCENT if item.get("highlight") else self.SECONDARY
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Pt(10), line_y - Pt(8), Pt(20), Pt(20))
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_color
            dot.line.fill.background()
            tb_left = MARGIN_L + int(item_w * i)
            yr = self._add_textbox(slide, tb_left, line_y + Pt(20), int(item_w), Inches(0.4))
            p = yr.text_frame.paragraphs[0]
            p.text = item.get("year", "")
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(13))
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY
            p.alignment = PP_ALIGN.CENTER
            txt = self._add_textbox(slide, tb_left, line_y + Pt(50), int(item_w), Inches(0.5))
            p2 = txt.text_frame.paragraphs[0]
            p2.text = item.get("text", "")
            p2.font.name = self.FONT
            p2.font.size = self._fs(Pt(11))
            p2.font.color.rgb = self.FG
            p2.alignment = PP_ALIGN.CENTER
            txt.text_frame.word_wrap = True
            if item.get("detail"):
                dtl = self._add_textbox(slide, tb_left, line_y + Pt(85), int(item_w), Inches(0.8))
                p3 = dtl.text_frame.paragraphs[0]
                p3.text = item["detail"]
                p3.font.name = self.FONT
                p3.font.size = self._fs(Pt(9))
                p3.font.color.rgb = self.MUTED
                p3.alignment = PP_ALIGN.CENTER
                dtl.text_frame.word_wrap = True

    def build_timeline_v(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.timeline_items)
        if n == 0:
            return
        line_x = MARGIN_L + Inches(0.15)
        row_h = min(Inches(1.2), BODY_H / max(n, 1))
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_x, BODY_TOP, Pt(3), int(row_h * n))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xDE, 0xE2, 0xE6)
        line.line.fill.background()
        for i, item in enumerate(sd.timeline_items):
            ry = BODY_TOP + int(row_h * i)
            dot_color = self.ACCENT if item.get("highlight") else self.SECONDARY
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, line_x - Pt(5), ry + Pt(6), Pt(12), Pt(12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_color
            dot.line.fill.background()
            content_left = MARGIN_L + Inches(0.6)
            content_w = CONTENT_W - Inches(0.6)
            yr = self._add_textbox(slide, content_left, ry, Inches(1.2), Inches(0.3))
            p = yr.text_frame.paragraphs[0]
            p.text = item.get("year", "")
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(13))
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY
            txt = self._add_textbox(slide, content_left + Inches(1.3), ry, content_w - Inches(1.3), Inches(0.3))
            p2 = txt.text_frame.paragraphs[0]
            p2.text = item.get("text", "")
            p2.font.name = self.FONT
            p2.font.size = self._fs(Pt(13))
            p2.font.color.rgb = self.FG
            txt.text_frame.word_wrap = True
            if item.get("detail"):
                dtl = self._add_textbox(slide, content_left + Inches(1.3), ry + Inches(0.3), content_w - Inches(1.3), Inches(0.4))
                p3 = dtl.text_frame.paragraphs[0]
                p3.text = item["detail"]
                p3.font.name = self.FONT
                p3.font.size = self._fs(Pt(10))
                p3.font.color.rgb = self.MUTED
                dtl.text_frame.word_wrap = True

    def build_end(self, sd: SlideData):
        slide = self._blank_slide()
        if self.LAYOUT.end_bg == "dark":
            self._set_bg(slide, self.PRIMARY)
        elif self.LAYOUT.end_bg == "light":
            self._set_bg(slide, self.LIGHT)
        is_dark = self.LAYOUT.end_bg == "dark"
        tb = self._add_textbox(slide, Inches(1), Inches(2), SW - Inches(2), Inches(2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = sd.h1 or "Thank You"
        p.font.name = self.FONT_HEAD
        p.font.size = self._fs(Pt(48))
        p.font.bold = True
        p.font.color.rgb = self.WHITE if is_dark else self.PRIMARY
        p.alignment = PP_ALIGN.CENTER
        remaining = []
        for line in sd.raw.split("\n"):
            s = line.strip()
            if not s.startswith("#") and s:
                remaining.append(strip_html(s))
        if remaining:
            tb2 = self._add_textbox(slide, Inches(1), Inches(4), SW - Inches(2), Inches(2))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            first = True
            for line in remaining:
                p2 = tf2.paragraphs[0] if first else tf2.add_paragraph()
                first = False
                p2.text = line
                p2.font.name = self.FONT
                p2.font.size = self._fs(Pt(18))
                p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC) if is_dark else self.MUTED
                p2.alignment = PP_ALIGN.CENTER

    def build_zone_flow(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.zone_flow_items)
        if n == 0:
            return
        gap = Inches(0.3)
        arrow_w = Inches(0.3)
        total_arrows = arrow_w * (n - 1)
        box_w = (CONTENT_W - total_arrows - gap * (n - 1)) / n
        box_h = Inches(3.5)
        for i, item in enumerate(sd.zone_flow_items):
            x = MARGIN_L + i * (box_w + gap + arrow_w)
            self._add_zone_box(slide, int(x), BODY_TOP, int(box_w), int(box_h),
                              label=item.get("label", ""), body=item.get("body", ""))
            if i < n - 1:
                ax = int(x + box_w + gap / 2)
                ay = int(BODY_TOP + box_h / 2 - Pt(10))
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, int(arrow_w), Pt(20))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.MUTED
                arrow.line.fill.background()
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_zone_compare(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        box_w = (CONTENT_W - Inches(1.2)) / 2
        box_h = Inches(4.0)
        self._add_zone_box(slide, MARGIN_L, BODY_TOP, int(box_w), int(box_h),
                          label=sd.zone_compare.get("left_label", ""),
                          body=sd.zone_compare.get("left_body", ""))
        vs_x = int(MARGIN_L + box_w + Inches(0.1))
        vs_tb = self._add_textbox(slide, vs_x, int(BODY_TOP + box_h / 2 - Pt(15)), Inches(1.0), Pt(30))
        p = vs_tb.text_frame.paragraphs[0]
        p.text = sd.zone_compare.get("vs_text", "VS")
        p.font.name = self.FONT_HEAD
        p.font.size = self._fs(Pt(18))
        p.font.bold = True
        p.font.color.rgb = self.ACCENT
        p.alignment = PP_ALIGN.CENTER
        self._add_zone_box(slide, int(MARGIN_L + box_w + Inches(1.2)), BODY_TOP, int(box_w), int(box_h),
                          label=sd.zone_compare.get("right_label", ""),
                          body=sd.zone_compare.get("right_body", ""))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_zone_matrix(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        cells = sd.zone_matrix.get("cells", [{}, {}, {}, {}])
        gap = Inches(0.2)
        cell_w = (CONTENT_W - gap) / 2
        cell_h = (BODY_H - gap) / 2
        positions = [
            (MARGIN_L, BODY_TOP),
            (int(MARGIN_L + cell_w + gap), BODY_TOP),
            (MARGIN_L, int(BODY_TOP + cell_h + gap)),
            (int(MARGIN_L + cell_w + gap), int(BODY_TOP + cell_h + gap)),
        ]
        for i, (x, y) in enumerate(positions):
            if i < len(cells):
                self._add_zone_box(slide, int(x), int(y), int(cell_w), int(cell_h),
                                  label=cells[i].get("label", ""),
                                  body=cells[i].get("body", ""))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_zone_process(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.zone_process_items)
        if n == 0:
            return
        gap = Inches(0.25)
        box_w = (CONTENT_W - gap * (n - 1)) / n
        box_h = Inches(4.0)
        for i, item in enumerate(sd.zone_process_items):
            x = MARGIN_L + i * (box_w + gap)
            num_h = Inches(0.6)
            num = self._add_textbox(slide, int(x), BODY_TOP, int(box_w), int(num_h))
            p = num.text_frame.paragraphs[0]
            p.text = item.get("step", str(i + 1))
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(24))
            p.font.bold = True
            p.font.color.rgb = self.ACCENT
            p.alignment = PP_ALIGN.CENTER
            self._add_zone_box(slide, int(x), int(BODY_TOP + num_h), int(box_w), int(box_h - num_h),
                              label=item.get("title", ""), body=item.get("body", ""))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_agenda(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.agenda_items:
            return
        tb = self._add_textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, BODY_H)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(sd.agenda_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(14)
            run_num = p.add_run()
            run_num.text = f"{i + 1}  "
            run_num.font.name = self.FONT_HEAD
            run_num.font.size = self._fs(Pt(24))
            run_num.font.bold = True
            run_num.font.color.rgb = self.SECONDARY
            run_txt = p.add_run()
            run_txt.text = item
            run_txt.font.name = self.FONT
            run_txt.font.size = self._fs(Pt(20))
            run_txt.font.color.rgb = self.FG

    def build_rq(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if sd.rq_main:
            tb = self._add_textbox(slide, Inches(1.5), Inches(2.5), SW - Inches(3), Inches(2))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd.rq_main
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(28))
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY
            p.alignment = PP_ALIGN.CENTER
        if sd.rq_sub:
            tb2 = self._add_textbox(slide, Inches(1.5), Inches(4.5), SW - Inches(3), Inches(1.5))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = sd.rq_sub
            p2.font.name = self.FONT
            p2.font.size = self._fs(Pt(16))
            p2.font.color.rgb = self.MUTED
            p2.alignment = PP_ALIGN.CENTER

    def build_result_dual(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.result_dual_items)
        if n == 0:
            return
        gap = Inches(0.4)
        col_w = (CONTENT_W - gap * (n - 1)) / n
        for i, item in enumerate(sd.result_dual_items):
            x = MARGIN_L + i * (col_w + gap)
            img_file = self._resolve_image(item.get("image", ""))
            if img_file:
                from PIL import Image
                with Image.open(img_file) as im:
                    iw, ih = im.size
                max_w = int(col_w * 0.95)
                max_h = int(Inches(3.5))
                scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
                pw = int(iw * scale * 914400 / 96)
                ph = int(ih * scale * 914400 / 96)
                img_left = int(x) + (int(col_w) - pw) // 2
                slide.shapes.add_picture(img_file, img_left, BODY_TOP, pw, ph)
            cap_text = item.get("caption", "")
            if cap_text:
                ctb = self._add_textbox(slide, int(x), BODY_TOP + Inches(3.8), int(col_w), Inches(0.5))
                cp = ctb.text_frame.paragraphs[0]
                cp.text = cap_text
                cp.font.name = self.FONT
                cp.font.size = self._fs(Pt(12))
                cp.font.color.rgb = self.FG
                cp.alignment = PP_ALIGN.CENTER

    def build_summary(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.summary_points:
            return
        tb = self._add_textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, BODY_H)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, pt in enumerate(sd.summary_points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(12)
            bdr = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                MARGIN_L, int(BODY_TOP + Inches(0.65) * i), Pt(5), Inches(0.5))
            bdr.fill.solid()
            bdr.fill.fore_color.rgb = self.ACCENT
            bdr.line.fill.background()
            run = p.add_run()
            run.text = f"  {pt}"
            run.font.name = self.FONT
            run.font.size = self._fs(Pt(18))
            run.font.color.rgb = self.FG

    def build_appendix(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1, color=self.MUTED)
        if sd.appendix_label:
            lbl = self._add_textbox(slide, CONTENT_W - Inches(1), MARGIN_T, Inches(2), Inches(0.3))
            p = lbl.text_frame.paragraphs[0]
            p.text = sd.appendix_label
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(10))
            p.font.color.rgb = self.MUTED
            p.alignment = PP_ALIGN.RIGHT
        if sd.table_rows:
            self.build_table(sd)
        elif sd.body_lines:
            self._add_body_text(slide, sd.body_lines, size=SZ_SMALL)

    def build_overview(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        cur_top = BODY_TOP
        if sd.overview_text:
            tb = self._add_textbox(slide, MARGIN_L, cur_top, CONTENT_W, Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd.overview_text
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(16))
            p.font.color.rgb = self.SECONDARY
            cur_top += Inches(0.8)
        left_w = CONTENT_W * 0.55
        right_w = CONTENT_W * 0.4
        if sd.image_path:
            img_file = self._resolve_image(sd.image_path)
            if img_file:
                from PIL import Image
                with Image.open(img_file) as im:
                    iw, ih = im.size
                max_w = int(left_w * 0.95)
                max_h = int(Inches(3.5))
                scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
                pw = int(iw * scale * 914400 / 96)
                ph = int(ih * scale * 914400 / 96)
                slide.shapes.add_picture(img_file, MARGIN_L, int(cur_top), pw, ph)
        if sd.overview_points:
            pts_left = int(MARGIN_L + left_w + Inches(0.3))
            tb = self._add_textbox(slide, pts_left, int(cur_top), int(right_w), BODY_H)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, pt in enumerate(sd.overview_points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(8)
                self._set_rich_text(p, f"\u2022 {pt}", Pt(14), self.FG)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_result(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        cur_top = BODY_TOP
        if sd.result_text:
            tb = self._add_textbox(slide, MARGIN_L, cur_top, CONTENT_W, Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd.result_text
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(16))
            p.font.color.rgb = self.SECONDARY
            cur_top += Inches(0.8)
        left_w = CONTENT_W * 0.55
        if sd.result_figure:
            img_file = self._resolve_image(sd.result_figure)
            if img_file:
                from PIL import Image
                with Image.open(img_file) as im:
                    iw, ih = im.size
                max_w = int(left_w * 0.95)
                max_h = int(Inches(3.5))
                scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
                pw = int(iw * scale * 914400 / 96)
                ph = int(ih * scale * 914400 / 96)
                slide.shapes.add_picture(img_file, MARGIN_L, int(cur_top), pw, ph)
        if sd.result_analysis:
            right_x = int(MARGIN_L + left_w + Inches(0.3))
            right_w = int(CONTENT_W * 0.4)
            tb = self._add_textbox(slide, right_x, int(cur_top), right_w, BODY_H)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, pt in enumerate(sd.result_analysis):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(8)
                self._set_rich_text(p, f"\u2022 {pt}", Pt(14), self.FG)
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_steps(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.steps_items)
        if n == 0:
            return
        gap = Inches(0.25)
        box_w = (CONTENT_W - gap * (n - 1)) / n
        box_h = Inches(4.0)
        for i, item in enumerate(sd.steps_items):
            x = MARGIN_L + i * (box_w + gap)
            num_tb = self._add_textbox(slide, int(x), BODY_TOP, int(box_w), Inches(0.5))
            p = num_tb.text_frame.paragraphs[0]
            p.text = item.get("num", str(i + 1))
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(28))
            p.font.bold = True
            p.font.color.rgb = self.ACCENT
            p.alignment = PP_ALIGN.CENTER
            self._add_zone_box(slide, int(x), int(BODY_TOP + Inches(0.6)), int(box_w), int(box_h - Inches(0.6)),
                              label=item.get("title", ""), body=item.get("body", ""))
        if sd.footnote:
            self._add_footnote(slide, sd.footnote)

    def build_quote(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        mark = self._add_textbox(slide, Inches(1.5), Inches(1.8), Inches(1), Inches(1))
        p = mark.text_frame.paragraphs[0]
        p.text = "\u201C"
        p.font.name = self.FONT_HEAD
        p.font.size = self._fs(Pt(72))
        p.font.color.rgb = self.SECONDARY
        if sd.quote_text:
            tb = self._add_textbox(slide, Inches(2.0), Inches(2.5), SW - Inches(4), Inches(3))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd.quote_text
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(22))
            p.font.color.rgb = self.FG
        if sd.quote_source:
            stb = self._add_textbox(slide, Inches(2.0), Inches(5.5), SW - Inches(4), Inches(0.5))
            p = stb.text_frame.paragraphs[0]
            p.text = f"\u2014 {sd.quote_source}"
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(14))
            p.font.color.rgb = self.MUTED
            p.alignment = PP_ALIGN.RIGHT

    def build_history(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.history_items)
        if n == 0:
            return
        row_h = min(Inches(0.7), BODY_H / max(n, 1))
        for i, item in enumerate(sd.history_items):
            y = BODY_TOP + int(row_h * i)
            yr = self._add_textbox(slide, MARGIN_L, y, Inches(1.5), int(row_h))
            p = yr.text_frame.paragraphs[0]
            p.text = item.get("year", "")
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(16))
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY
            p.alignment = PP_ALIGN.RIGHT
            ev = self._add_textbox(slide, MARGIN_L + Inches(1.8), y, CONTENT_W - Inches(1.8), int(row_h))
            p2 = ev.text_frame.paragraphs[0]
            p2.text = item.get("event", "")
            p2.font.name = self.FONT
            p2.font.size = self._fs(Pt(14))
            p2.font.color.rgb = self.FG
            ev.text_frame.word_wrap = True

    def build_panorama(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        img_file = self._resolve_image(sd.image_path) if sd.image_path else None
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(CONTENT_W)
            max_h = int(Inches(4.5))
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96)
            ph = int(ih * scale * 914400 / 96)
            left = (SW - pw) // 2
            slide.shapes.add_picture(img_file, left, BODY_TOP, pw, ph)
        if sd.panorama_text:
            tb = self._add_textbox(slide, MARGIN_L, SH - Inches(1.2), CONTENT_W, Inches(0.6))
            p = tb.text_frame.paragraphs[0]
            p.text = sd.panorama_text
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(14))
            p.font.color.rgb = self.MUTED
            p.alignment = PP_ALIGN.CENTER

    def build_kpi(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.kpi_items)
        if n == 0:
            return
        gap = Inches(0.4)
        box_w = (CONTENT_W - gap * (n - 1)) / n
        box_h = Inches(3.0)
        kpi_top = int(BODY_TOP + (BODY_H - box_h) / 2)
        for i, item in enumerate(sd.kpi_items):
            x = MARGIN_L + i * (box_w + gap)
            self._add_zone_box(slide, int(x), kpi_top, int(box_w), int(box_h), label="", body="")
            vtb = self._add_textbox(slide, int(x), kpi_top + Inches(0.5), int(box_w), Inches(1.2))
            p = vtb.text_frame.paragraphs[0]
            p.text = item.get("value", "")
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(36))
            p.font.bold = True
            p.font.color.rgb = self.ACCENT
            p.alignment = PP_ALIGN.CENTER
            ltb = self._add_textbox(slide, int(x), kpi_top + Inches(1.8), int(box_w), Inches(0.6))
            p2 = ltb.text_frame.paragraphs[0]
            p2.text = item.get("label", "")
            p2.font.name = self.FONT
            p2.font.size = self._fs(Pt(14))
            p2.font.color.rgb = self.MUTED
            p2.alignment = PP_ALIGN.CENTER

    def build_pros_cons(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        half_w = (CONTENT_W - Inches(0.4)) / 2
        box_h = Inches(4.5)
        for side, items, color, label_text in [
            ("left", sd.pros_items, RGBColor(0x2E, 0x7D, 0x32), "Pros"),
            ("right", sd.cons_items, RGBColor(0xC6, 0x28, 0x28), "Cons"),
        ]:
            x = MARGIN_L if side == "left" else int(MARGIN_L + half_w + Inches(0.4))
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), BODY_TOP, int(half_w), int(box_h))
            bg.adjustments[0] = 0.02
            bg.fill.background()
            bg.line.color.rgb = color
            bg.line.width = Pt(1.5)
            lbl = self._add_textbox(slide, int(x) + Pt(16), BODY_TOP + Pt(10), int(half_w) - Pt(32), Inches(0.4))
            p = lbl.text_frame.paragraphs[0]
            p.text = label_text
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(18))
            p.font.bold = True
            p.font.color.rgb = color
            itb = self._add_textbox(slide, int(x) + Pt(16), BODY_TOP + Inches(0.6), int(half_w) - Pt(32), int(box_h) - Inches(0.8))
            tf = itb.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"\u2022 {item}"
                p.font.name = self.FONT
                p.font.size = self._fs(Pt(14))
                p.font.color.rgb = self.FG
                p.space_before = Pt(6)

    def build_definition(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if sd.def_term:
            tb = self._add_textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, Inches(0.6))
            p = tb.text_frame.paragraphs[0]
            p.text = sd.def_term
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(28))
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY
        if sd.def_body:
            tb2 = self._add_textbox(slide, MARGIN_L + Inches(0.3), BODY_TOP + Inches(0.9), CONTENT_W - Inches(0.3), Inches(3))
            tf2 = tb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = sd.def_body
            p2.font.name = self.FONT
            p2.font.size = self._fs(Pt(18))
            p2.font.color.rgb = self.FG
        if sd.def_note:
            ntb = self._add_textbox(slide, MARGIN_L, SH - Inches(1.5), CONTENT_W, Inches(0.8))
            p = ntb.text_frame.paragraphs[0]
            p.text = sd.def_note
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(12))
            p.font.color.rgb = self.MUTED
            ntb.text_frame.word_wrap = True

    def build_diagram(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        img_file = self._resolve_image(sd.image_path) if sd.image_path else None
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(CONTENT_W * 0.9)
            max_h = int(Inches(4.5))
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96)
            ph = int(ih * scale * 914400 / 96)
            left = (SW - pw) // 2
            slide.shapes.add_picture(img_file, left, BODY_TOP, pw, ph)
        if sd.caption:
            ctb = self._add_textbox(slide, MARGIN_L, SH - Inches(1.0), CONTENT_W, Inches(0.5))
            p = ctb.text_frame.paragraphs[0]
            p.text = sd.caption
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(12))
            p.font.color.rgb = self.FG
            p.alignment = PP_ALIGN.CENTER

    def build_gallery_img(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.gallery_items)
        if n == 0:
            return
        cols = min(n, 3)
        rows_n = (n + cols - 1) // cols
        gap = Inches(0.3)
        cell_w = (CONTENT_W - gap * (cols - 1)) / cols
        cell_h = (BODY_H - gap * (rows_n - 1)) / rows_n
        for idx, item in enumerate(sd.gallery_items):
            r, c = divmod(idx, cols)
            x = MARGIN_L + c * (cell_w + gap)
            y = BODY_TOP + r * (cell_h + gap)
            img_file = self._resolve_image(item.get("image", ""))
            if img_file:
                from PIL import Image
                with Image.open(img_file) as im:
                    iw, ih = im.size
                max_w = int(cell_w * 0.95)
                max_h = int(cell_h * 0.8)
                scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
                pw = int(iw * scale * 914400 / 96)
                ph = int(ih * scale * 914400 / 96)
                slide.shapes.add_picture(img_file, int(x), int(y), pw, ph)

    def build_highlight(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if sd.highlight_text:
            tb = self._add_textbox(slide, Inches(1.5), Inches(2.5), SW - Inches(3), Inches(3))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd.highlight_text
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(32))
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY
            p.alignment = PP_ALIGN.CENTER

    def build_checklist(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if not sd.checklist_items:
            return
        tb = self._add_textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, BODY_H)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(sd.checklist_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            check = "\u2611" if item.get("done") else "\u2610"
            p.text = f"{check}  {item.get('text', '')}"
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(16))
            p.font.color.rgb = self.MUTED if item.get("done") else self.FG
            p.space_before = Pt(8)

    def build_annotation(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        fig_w = CONTENT_W * 0.55
        img_file = self._resolve_image(sd.annotation_figure) if sd.annotation_figure else None
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(fig_w * 0.95)
            max_h = int(Inches(4.5))
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96)
            ph = int(ih * scale * 914400 / 96)
            slide.shapes.add_picture(img_file, MARGIN_L, BODY_TOP, pw, ph)
        if sd.annotation_notes:
            notes_x = int(MARGIN_L + fig_w + Inches(0.3))
            notes_w = int(CONTENT_W - fig_w - Inches(0.3))
            tb = self._add_textbox(slide, notes_x, BODY_TOP, notes_w, BODY_H)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, note in enumerate(sd.annotation_notes):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"\u2022 {note}"
                p.font.name = self.FONT
                p.font.size = self._fs(Pt(13))
                p.font.color.rgb = self.FG
                p.space_before = Pt(8)

    def build_before_after(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        half_w = (CONTENT_W - Inches(0.8)) / 2
        box_h = Inches(4.5)
        for side, data, label_default in [
            ("left", sd.ba_before, "Before"),
            ("right", sd.ba_after, "After"),
        ]:
            x = MARGIN_L if side == "left" else int(MARGIN_L + half_w + Inches(0.8))
            label = data.get("label", label_default) if data else label_default
            body = data.get("body", "") if data else ""
            self._add_zone_box(slide, int(x), BODY_TOP, int(half_w), int(box_h), label=label, body=body)
        arrow_x = int(MARGIN_L + half_w + Inches(0.1))
        arrow_y = int(BODY_TOP + box_h / 2 - Pt(15))
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.6), Pt(30))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = self.ACCENT
        arrow.line.fill.background()

    def build_funnel(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.funnel_items)
        if n == 0:
            return
        max_w = CONTENT_W
        min_w = max_w * 0.3
        row_h = min(Inches(0.9), BODY_H / max(n, 1))
        for i, item in enumerate(sd.funnel_items):
            frac = 1 - (i / max(n - 1, 1)) * 0.7
            w = int(max_w * frac)
            x = int(MARGIN_L + (max_w - w) / 2)
            y = int(BODY_TOP + row_h * i)
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, int(row_h * 0.85))
            bg.adjustments[0] = 0.05
            bg.fill.solid()
            t = i / max(n - 1, 1)
            r = int(self.PRIMARY.red * (1 - t) + self.ACCENT.red * t) if hasattr(self.PRIMARY, 'red') else 0x16
            g = int(self.PRIMARY.green * (1 - t) + self.ACCENT.green * t) if hasattr(self.PRIMARY, 'green') else 0x21
            b = int(self.PRIMARY.blue * (1 - t) + self.ACCENT.blue * t) if hasattr(self.PRIMARY, 'blue') else 0x3e
            bg.fill.fore_color.rgb = RGBColor(min(r, 255), min(g, 255), min(b, 255))
            bg.line.fill.background()
            tb = self._add_textbox(slide, x + Pt(16), y + Pt(4), w - Pt(32), int(row_h * 0.85) - Pt(8))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            run_l = tf.paragraphs[0].add_run()
            run_l.text = item.get("label", "")
            run_l.font.name = self.FONT_HEAD
            run_l.font.size = self._fs(Pt(14))
            run_l.font.bold = True
            run_l.font.color.rgb = self.WHITE
            if item.get("value"):
                run_v = tf.paragraphs[0].add_run()
                run_v.text = f"  {item['value']}"
                run_v.font.name = self.FONT
                run_v.font.size = self._fs(Pt(12))
                run_v.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

    def build_stack(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.stack_items)
        if n == 0:
            return
        gap = Inches(0.1)
        row_h = min(Inches(1.0), (BODY_H - gap * (n - 1)) / n)
        for i, item in enumerate(sd.stack_items):
            y = int(BODY_TOP + (row_h + gap) * (n - 1 - i))
            self._add_zone_box(slide, MARGIN_L, y, CONTENT_W, int(row_h),
                              label=item.get("name", ""), body=item.get("desc", ""))

    def build_card_grid(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.card_items)
        if n == 0:
            return
        cols = 2 if n <= 4 else 3
        rows_n = (n + cols - 1) // cols
        gap = Inches(0.25)
        card_w = (CONTENT_W - gap * (cols - 1)) / cols
        card_h = (BODY_H - gap * (rows_n - 1)) / rows_n
        for idx, item in enumerate(sd.card_items):
            r, c = divmod(idx, cols)
            x = MARGIN_L + c * (card_w + gap)
            y = BODY_TOP + r * (card_h + gap)
            self._add_zone_box(slide, int(x), int(y), int(card_w), int(card_h),
                              label=item.get("title", ""), body=item.get("body", ""))

    def build_split_text(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        half_w = (CONTENT_W - Inches(0.3)) / 2
        for side, data in [("left", sd.split_left), ("right", sd.split_right)]:
            x = MARGIN_L if side == "left" else int(MARGIN_L + half_w + Inches(0.3))
            label = data.get("label", "") if data else ""
            body = data.get("body", "") if data else ""
            self._add_zone_box(slide, int(x), BODY_TOP, int(half_w), BODY_H, label=label, body=body)

    def build_code(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        code_h = Inches(4.0) if sd.code_desc else Inches(5.0)
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L, BODY_TOP, CONTENT_W, int(code_h))
        bg.adjustments[0] = 0.01
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
        bg.line.fill.background()
        tb = self._add_textbox(slide, MARGIN_L + Pt(20), BODY_TOP + Pt(16), CONTENT_W - Pt(40), int(code_h) - Pt(32))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sd.code_text
        p.font.name = self.FONT_MONO
        p.font.size = self._fs(Pt(12))
        p.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)
        if sd.code_desc:
            dtb = self._add_textbox(slide, MARGIN_L, int(BODY_TOP + code_h + Inches(0.2)), CONTENT_W, Inches(0.8))
            p2 = dtb.text_frame.paragraphs[0]
            p2.text = sd.code_desc
            p2.font.name = self.FONT
            p2.font.size = self._fs(Pt(13))
            p2.font.color.rgb = self.MUTED
            dtb.text_frame.word_wrap = True

    def build_multi_result(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        n = len(sd.multi_result_items)
        if n == 0:
            return
        gap = Inches(0.3)
        box_w = (CONTENT_W - gap * (n - 1)) / n
        box_h = Inches(4.0)
        for i, item in enumerate(sd.multi_result_items):
            x = MARGIN_L + i * (box_w + gap)
            self._add_zone_box(slide, int(x), BODY_TOP, int(box_w), int(box_h), label="", body="")
            mtb = self._add_textbox(slide, int(x), BODY_TOP + Inches(0.3), int(box_w), Inches(0.5))
            p = mtb.text_frame.paragraphs[0]
            p.text = item.get("metric", "")
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(12))
            p.font.color.rgb = self.MUTED
            p.alignment = PP_ALIGN.CENTER
            vtb = self._add_textbox(slide, int(x), BODY_TOP + Inches(0.9), int(box_w), Inches(1.0))
            p2 = vtb.text_frame.paragraphs[0]
            p2.text = item.get("value", "")
            p2.font.name = self.FONT_HEAD
            p2.font.size = self._fs(Pt(32))
            p2.font.bold = True
            p2.font.color.rgb = self.ACCENT
            p2.alignment = PP_ALIGN.CENTER
            dtb = self._add_textbox(slide, int(x) + Pt(10), BODY_TOP + Inches(2.2), int(box_w) - Pt(20), Inches(1.5))
            p3 = dtb.text_frame.paragraphs[0]
            p3.text = item.get("desc", "")
            p3.font.name = self.FONT
            p3.font.size = self._fs(Pt(12))
            p3.font.color.rgb = self.FG
            p3.alignment = PP_ALIGN.CENTER
            dtb.text_frame.word_wrap = True

    def build_takeaway(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        if sd.takeaway_main:
            tb = self._add_textbox(slide, Inches(1.5), Inches(2.0), SW - Inches(3), Inches(1.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = sd.takeaway_main
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(28))
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY
            p.alignment = PP_ALIGN.CENTER
        if sd.takeaway_points:
            ptb = self._add_textbox(slide, Inches(2.0), Inches(3.8), SW - Inches(4), Inches(3))
            tf = ptb.text_frame
            tf.word_wrap = True
            for i, pt in enumerate(sd.takeaway_points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"\u2022 {pt}"
                p.font.name = self.FONT
                p.font.size = self._fs(Pt(16))
                p.font.color.rgb = self.FG
                p.space_before = Pt(8)

    def build_profile(self, sd: SlideData):
        slide = self._blank_slide()
        if sd.h1:
            self._add_title(slide, sd.h1)
        left_w = Inches(3.5)
        right_x = MARGIN_L + left_w + Inches(0.5)
        right_w = CONTENT_W - left_w - Inches(0.5)
        img_file = self._resolve_image(sd.image_path) if sd.image_path else None
        if img_file:
            from PIL import Image
            with Image.open(img_file) as im:
                iw, ih = im.size
            max_w = int(left_w * 0.9)
            max_h = int(Inches(3.5))
            scale = min(max_w / (iw * 914400 / 96), max_h / (ih * 914400 / 96))
            pw = int(iw * scale * 914400 / 96)
            ph = int(ih * scale * 914400 / 96)
            slide.shapes.add_picture(img_file, MARGIN_L, BODY_TOP, pw, ph)
        cur_y = BODY_TOP
        if sd.profile_name:
            ntb = self._add_textbox(slide, int(right_x), int(cur_y), int(right_w), Inches(0.5))
            p = ntb.text_frame.paragraphs[0]
            p.text = sd.profile_name
            p.font.name = self.FONT_HEAD
            p.font.size = self._fs(Pt(24))
            p.font.bold = True
            p.font.color.rgb = self.PRIMARY
            cur_y += Inches(0.6)
        if sd.profile_affiliation:
            atb = self._add_textbox(slide, int(right_x), int(cur_y), int(right_w), Inches(0.4))
            p = atb.text_frame.paragraphs[0]
            p.text = sd.profile_affiliation
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(14))
            p.font.color.rgb = self.MUTED
            cur_y += Inches(0.6)
        if sd.profile_bio:
            btb = self._add_textbox(slide, int(right_x), int(cur_y), int(right_w), int(BODY_H - (cur_y - BODY_TOP)))
            btf = btb.text_frame
            btf.word_wrap = True
            first = True
            for item in sd.profile_bio:
                if first:
                    p = btf.paragraphs[0]
                    first = False
                else:
                    p = btf.add_paragraph()
                run = p.add_run()
                run.text = "\u2022 " + item
                run.font.name = self.FONT
                run.font.size = self._fs(Pt(15))
                run.font.color.rgb = self.FG
                p.space_before = Pt(6)

    # ══════════════════════════════════════════════
    # Build all slides
    # ══════════════════════════════════════════════
    BUILDERS = {
        "title": "build_title",
        "divider": "build_divider",
        "cols-2": "build_columns",
        "cols-2-wide-l": "build_columns",
        "cols-2-wide-r": "build_columns",
        "cols-3": "build_columns",
        "sandwich": "build_sandwich",
        "equation": "build_equation",
        "equations": "build_equations",
        "figure": "build_figure",
        "table-slide": "build_table",
        "references": "build_references",
        "timeline-h": "build_timeline_h",
        "timeline": "build_timeline_v",
        "end": "build_end",
        "zone-flow": "build_zone_flow",
        "zone-compare": "build_zone_compare",
        "zone-matrix": "build_zone_matrix",
        "zone-process": "build_zone_process",
        "agenda": "build_agenda",
        "rq": "build_rq",
        "result-dual": "build_result_dual",
        "summary": "build_summary",
        "appendix": "build_appendix",
        "overview": "build_overview",
        "result": "build_result",
        "steps": "build_steps",
        "quote": "build_quote",
        "history": "build_history",
        "panorama": "build_panorama",
        "kpi": "build_kpi",
        "pros-cons": "build_pros_cons",
        "definition": "build_definition",
        "diagram": "build_diagram",
        "gallery-img": "build_gallery_img",
        "highlight": "build_highlight",
        "checklist": "build_checklist",
        "annotation": "build_annotation",
        "before-after": "build_before_after",
        "funnel": "build_funnel",
        "stack": "build_stack",
        "card-grid": "build_card_grid",
        "split-text": "build_split_text",
        "code": "build_code",
        "multi-result": "build_multi_result",
        "takeaway": "build_takeaway",
        "profile": "build_profile",
    }

    def build_all(self, slides: list[SlideData]):
        for sd in slides:
            method_name = self.BUILDERS.get(sd.slide_class, "build_default")
            getattr(self, method_name)(sd)
        self._add_global_footer()

    def _add_global_footer(self):
        for i, slide in enumerate(self.prs.slides):
            if i == 0 or i == len(self.prs.slides) - 1:
                continue
            tb = self._add_textbox(slide, int(MARGIN_L), int(SH - Inches(0.4)), int(CONTENT_W), Inches(0.25))
            p = tb.text_frame.paragraphs[0]
            p.text = f"{i + 1}"
            p.font.name = self.FONT
            p.font.size = self._fs(Pt(8))
            p.font.color.rgb = self.MUTED
            p.alignment = PP_ALIGN.RIGHT
