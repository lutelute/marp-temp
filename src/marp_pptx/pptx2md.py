"""PPTX → Markdown (best-effort) for learning-data bootstrapping.

This is NOT a round-trip of the builder's MD→PPTX path. Generic PPTX
files don't carry semantic type hints, so we produce best-effort
Markdown using text extraction + heuristic type inference, which the
user can then edit (flipping <!-- _class: ... --> directives) to build
a labelled (PPTX, MD) pair for training downstream models.
"""
from __future__ import annotations

import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


@dataclass
class ExtractedSlide:
    index: int
    title: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    numbered: list[str] = field(default_factory=list)
    paragraphs: list[str] = field(default_factory=list)
    table_rows: list[list[str]] = field(default_factory=list)
    images: list[str] = field(default_factory=list)
    # Raw shape dump for heuristic analysis
    shapes: list[dict] = field(default_factory=list)
    # Detected type (heuristic)
    inferred_class: str | None = None


# ── Heuristic type detection ──

_NUMBERIC_RE = re.compile(r"^[\-+]?\d+[.,]?\d*\s*(%|x|倍|ms|s|ms|GB|MB|KB|人|件)?$")
_PROS_HINTS = {"pros", "長所", "利点", "メリット", "◯", "○"}
_CONS_HINTS = {"cons", "短所", "欠点", "デメリット", "×", "✗"}
_BEFORE_HINTS = {"before", "改善前", "従来", "現状"}
_AFTER_HINTS = {"after", "改善後", "提案", "新"}


def _extract_text_from_shape(shape) -> str:
    if not shape.has_text_frame:
        return ""
    parts = []
    for p in shape.text_frame.paragraphs:
        parts.append(p.text)
    return "\n".join(parts).strip()


def _is_bullet(para) -> bool:
    # A paragraph with bullet formatting will have <a:buChar> in pPr
    try:
        pPr = para._pPr
        if pPr is None:
            return False
        from pptx.oxml.ns import qn
        return pPr.find(qn("a:buChar")) is not None or pPr.find(qn("a:buAutoNum")) is not None
    except Exception:
        return False


def extract_slide(idx: int, slide) -> ExtractedSlide:
    es = ExtractedSlide(index=idx)
    largest_by_font = (None, 0.0)  # (text, font_size)
    second_largest = (None, 0.0)

    for shape in slide.shapes:
        entry = {
            "type": str(shape.shape_type),
            "name": shape.name,
            "left": shape.left or 0,
            "top": shape.top or 0,
            "width": shape.width or 0,
            "height": shape.height or 0,
            "text": "",
            "bullets": [],
            "numbered": [],
        }

        # Pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            entry["kind"] = "picture"
            try:
                img = shape.image
                ext = img.ext
                entry["image_ext"] = ext
                # Save later; for now just record placeholder
                es.images.append(f"image_{idx}_{shape.shape_id}.{ext}")
            except Exception:
                pass
            es.shapes.append(entry)
            continue

        # Tables
        if shape.has_table:
            entry["kind"] = "table"
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            entry["table"] = rows
            es.table_rows = rows  # keep just one table per slide (simplification)
            es.shapes.append(entry)
            continue

        # Text frames
        if shape.has_text_frame:
            text_all = _extract_text_from_shape(shape)
            entry["text"] = text_all
            if not text_all:
                es.shapes.append(entry)
                continue
            # Track largest (heuristic H1 / H2 detection via font size)
            for para in shape.text_frame.paragraphs:
                size = None
                for run in para.runs:
                    if run.font.size:
                        size = run.font.size.pt
                        break
                if size is None and para.font.size:
                    size = para.font.size.pt
                if size is None:
                    continue
                if size > largest_by_font[1]:
                    second_largest = largest_by_font
                    largest_by_font = (para.text.strip(), size)
                elif size > second_largest[1]:
                    second_largest = (para.text.strip(), size)

            # Per-paragraph classification
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if not t:
                    continue
                if _is_bullet(para):
                    entry["bullets"].append(t)
                    es.bullets.append(t)
                elif re.match(r"^\d+[.)]\s+", t):
                    entry["numbered"].append(t)
                    es.numbered.append(re.sub(r"^\d+[.)]\s+", "", t))
                else:
                    es.paragraphs.append(t)
            es.shapes.append(entry)

    # Choose title as largest-font text if it's reasonably short
    if largest_by_font[0] and len(largest_by_font[0]) < 80:
        es.title = largest_by_font[0]
        if second_largest[0] and second_largest[0] != es.title and len(second_largest[0]) < 120:
            es.subtitle = second_largest[0]

    # Remove title/subtitle from paragraphs list to avoid duplication
    if es.title:
        es.paragraphs = [p for p in es.paragraphs if p != es.title]
    if es.subtitle:
        es.paragraphs = [p for p in es.paragraphs if p != es.subtitle]

    return es


def infer_slide_type(es: ExtractedSlide, total_slides: int) -> str | None:
    """Heuristic mapping: extracted features → most likely slide class."""
    # Meta types by position
    if es.index == 0:
        return "title"
    if es.index == total_slides - 1:
        if es.title and len(es.title) < 30 and not es.bullets and not es.paragraphs:
            return "end"
    if es.table_rows:
        return "table-slide"
    if len(es.images) >= 2 and not es.paragraphs:
        return "gallery-img"
    if len(es.images) == 1 and (es.paragraphs or es.bullets):
        # Figure or diagram
        if any("図" in p or "fig" in p.lower() for p in es.paragraphs + es.bullets):
            return "figure"
        return "figure"
    if len(es.images) == 1:
        return "figure"

    # Check bullets for pros/cons pattern
    text_blob = " ".join(es.bullets + es.paragraphs).lower()
    if any(h in text_blob for h in _PROS_HINTS) and any(h in text_blob for h in _CONS_HINTS):
        return "pros-cons"

    # Before/After
    if any(h in text_blob for h in _BEFORE_HINTS) and any(h in text_blob for h in _AFTER_HINTS):
        return "before-after"

    # KPI: 3-6 large short numeric shapes
    numeric_shapes = []
    for sh in es.shapes:
        if sh.get("kind") not in (None, "text"):
            continue
        text = sh.get("text", "").strip()
        # Single short line, looks like a number
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        if 1 <= len(lines) <= 2 and any(_NUMBERIC_RE.match(l) for l in lines):
            numeric_shapes.append(sh)
    if 3 <= len(numeric_shapes) <= 6:
        return "kpi"

    # Agenda: numbered list of 3-7 items
    if 3 <= len(es.numbered) <= 10 and not es.bullets and not es.images:
        return "agenda"

    # Summary: numbered 2-5 items at end of deck
    if es.index >= total_slides * 0.6 and es.numbered and len(es.numbered) <= 7:
        return "summary"

    # Divider: very few shapes, short title, center
    n_text_shapes = sum(1 for s in es.shapes if s.get("text"))
    if n_text_shapes <= 2 and es.title and not es.bullets and not es.paragraphs:
        return "divider"

    # Quote: has long paragraph + short source line
    if len(es.paragraphs) == 2 and len(es.paragraphs[0]) > 30 and len(es.paragraphs[1]) < 40:
        if any(c in es.paragraphs[0] for c in "「\"'"):
            return "quote"

    # Equation: looks like LaTeX math
    if any("$$" in p or re.search(r"[\\\\^_{}]", p) for p in es.paragraphs):
        return "equation"

    # Default: fall through
    return None


def _escape_md(s: str) -> str:
    return s


def slide_to_md(es: ExtractedSlide, cls: str | None) -> str:
    """Convert an extracted slide to best-effort Markdown."""
    parts = []
    if cls:
        parts.append(f"<!-- _class: {cls} -->")
    if es.title:
        parts.append(f"# {es.title}")
    if es.subtitle and cls in (None, "title", "divider"):
        parts.append(f"## {es.subtitle}")

    # Type-specific rendering
    if cls == "table-slide" and es.table_rows:
        if es.table_rows:
            header = "| " + " | ".join(es.table_rows[0]) + " |"
            sep = "| " + " | ".join(["---"] * len(es.table_rows[0])) + " |"
            body_rows = ["| " + " | ".join(r) + " |" for r in es.table_rows[1:]]
            parts.append(header)
            parts.append(sep)
            parts.extend(body_rows)
    elif cls == "kpi":
        # Try to pair numeric shapes with label shapes
        items = []
        for sh in es.shapes:
            text = sh.get("text", "")
            lines = [l.strip() for l in text.split("\n") if l.strip()]
            if len(lines) == 1 and _NUMBERIC_RE.match(lines[0]):
                items.append({"value": lines[0], "label": ""})
            elif len(lines) == 2 and _NUMBERIC_RE.match(lines[0]):
                items.append({"value": lines[0], "label": lines[1]})
        if items:
            parts.append('<div class="kpi-container">')
            for it in items[:6]:
                parts.append(f'<div><span class="kpi-value">{it["value"]}</span>'
                             f'<span class="kpi-label">{it["label"]}</span></div>')
            parts.append('</div>')
    elif cls == "agenda" and es.numbered:
        parts.append('<div class="agenda-list">')
        for i, item in enumerate(es.numbered, 1):
            parts.append(f"{i}. {item}")
        parts.append('</div>')
    elif cls == "summary" and es.numbered:
        parts.append('<ol class="summary-points">')
        for item in es.numbered:
            parts.append(f"<li>{item}</li>")
        parts.append('</ol>')
    elif cls == "figure" and es.images:
        parts.append(f"![]({es.images[0]})")
        if es.paragraphs:
            parts.append(f'<div class="caption">{es.paragraphs[0]}</div>')
    elif cls == "gallery-img" and es.images:
        parts.append('<div class="gi-container">')
        for img in es.images[:6]:
            parts.append(f'<div>\n![]({img})\n</div>')
        parts.append('</div>')
    else:
        # Default: paragraphs + bullets + remaining
        for p in es.paragraphs:
            parts.append(p)
        for b in es.bullets:
            parts.append(f"- {b}")
        if es.numbered and cls not in ("agenda", "summary"):
            for i, n in enumerate(es.numbered, 1):
                parts.append(f"{i}. {n}")
        if es.images and cls not in ("figure", "gallery-img"):
            for img in es.images:
                parts.append(f"![]({img})")

    return "\n".join(parts)


def pptx_to_md(pptx_path: Path, extract_images_to: Path | None = None) -> str:
    """Convert a PPTX file to best-effort Markdown for our 49-type system.

    If extract_images_to is provided, embedded images are saved there
    (with safe filenames) and referenced as assets/<name>.
    """
    prs = Presentation(str(pptx_path))
    slides_extracted = [extract_slide(i, s) for i, s in enumerate(prs.slides)]
    total = len(slides_extracted)

    # Extract images if dir provided
    if extract_images_to is not None:
        extract_images_to.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(str(pptx_path)) as z:
            img_files = sorted(n for n in z.namelist() if n.startswith("ppt/media/"))
            for name in img_files:
                data = z.read(name)
                out = extract_images_to / Path(name).name
                out.write_bytes(data)

    # Build MD
    md_parts = ["---", "marp: true", "theme: academic", "math: katex", "---", ""]
    for es in slides_extracted:
        cls = infer_slide_type(es, total)
        es.inferred_class = cls
        md_parts.append(slide_to_md(es, cls))
        md_parts.append("")
        md_parts.append("---")
        md_parts.append("")
    # Remove trailing ---
    if md_parts and md_parts[-2] == "---":
        md_parts = md_parts[:-2]

    return "\n".join(md_parts)


def pptx_to_md_with_report(pptx_path: Path, extract_images_to: Path | None = None) -> dict:
    """Same as pptx_to_md, but returns a dict with metadata for training-data use."""
    md = pptx_to_md(pptx_path, extract_images_to=extract_images_to)
    prs = Presentation(str(pptx_path))
    slides_extracted = [extract_slide(i, s) for i, s in enumerate(prs.slides)]
    report = {
        "markdown": md,
        "slides": [
            {
                "index": es.index,
                "title": es.title,
                "inferred_class": infer_slide_type(es, len(slides_extracted)),
                "bullets": len(es.bullets),
                "numbered": len(es.numbered),
                "paragraphs": len(es.paragraphs),
                "images": len(es.images),
                "table_rows": len(es.table_rows),
            }
            for es in slides_extracted
        ],
    }
    return report
